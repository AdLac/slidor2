import openai
import time
from pptx import Presentation
import streamlit as st

# Clé API OpenAI
openai.api_key = st.secrets["OPENAI_API_KEY"]

def main():
    # Titre de l'application
    st.title('Générateur de Slides :sunglasses:')

    # Charger le modèle de présentation
    try:
        prs = Presentation('template_nexus.pptx')
    except Exception as e:
        st.error(f"Erreur lors du chargement du modèle de présentation : {e}")
        return

    # Définir la mise en page des diapositives
    slide_layout = prs.slide_layouts[12]

    # Section pour l'entrée des mots-clés
    st.subheader('Générer les diapositives : 1 ligne = 1 titre de diapositive')
    contexte = st.text_area("Contexte (aide à obtenir des résultats plus précis, par exemple : 'Vous travaillez pour ce client...')")
    keywords = st.text_area('Mots-clés (1 ligne, 1 titre)')

    if st.button('Générer les diapositives') and keywords:
        keywords = [kw.strip() for kw in keywords.split("\n") if kw.strip()]
        if not keywords:
            st.warning("Veuillez entrer au moins un mot-clé valide.")
            return

        counter = 0
        for keyword in keywords:
            counter += 1

            # Créer le prompt pour l'API OpenAI
            prompt = (
                f"{contexte}\n\n"
                f"Rédige un titre ainsi qu'un commentaire complet et détaillé à partir de cette idée : \"{keyword}\". "
                "Utilise le format suivant pour le titre : <T>titre-généré-ici</T> et pour le commentaire : <C>commentaire-généré-ici</C>."
            )

            # Appel à l'API OpenAI
            try:
                response = openai.ChatCompletion.create(
                    model="gpt-3.5-turbo",
                    messages=[
                        {"role": "system", "content": "Vous êtes un assistant utile."},
                        {"role": "user", "content": prompt}
                    ],
                    temperature=0.7,
                    max_tokens=500,
                    top_p=1,
                    frequency_penalty=0,
                    presence_penalty=0
                )
            except Exception as e:
                st.error(f"Erreur lors de l'appel à l'API OpenAI pour le mot-clé '{keyword}' : {e}")
                continue

            content = response.choices[0].message['content']

            # Extraire le titre et le commentaire
            try:
                title = content.split("<T>")[1].split("</T>")[0].strip()
                body = content.split("<C>")[1].split("</C>")[0].strip()
            except IndexError:
                st.error(f"Le format de la réponse pour le mot-clé '{keyword}' est incorrect. Assurez-vous que le modèle renvoie le titre et le commentaire dans les balises appropriées.")
                continue

            # Ajouter une diapositive à la présentation
            slide = prs.slides.add_slide(slide_layout)
            try:
                slide.placeholders[0].text = title  # Titre de la diapositive
                slide.placeholders[11].text = title  # Sous-titre de la diapositive
                slide.placeholders[21].text = body   # Contenu de la diapositive
            except IndexError as e:
                st.error(f"Erreur lors de l'ajout de contenu à la diapositive pour le mot-clé '{keyword}' : {e}")
                continue

            # Afficher la progression
            st.text(f"{counter} diapositive(s) générée(s).")

            # Pause pour éviter de surcharger l'API
            time.sleep(0.5)

        # Enregistrer la présentation modifiée
        output_file = 'template_nexus_modified.pptx'
        try:
            prs.save(output_file)
            st.success(f"La présentation modifiée a été enregistrée sous '{output_file}' !")
        except Exception as e:
            st.error(f"Erreur lors de l'enregistrement de la présentation : {e}")
            return

        # Bouton de téléchargement
        with open(output_file, "rb") as file:
            st.download_button(
                label="Télécharger la présentation",
                data=file,
                file_name=output_file,
                mime='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            )

if __name__ == '__main__':
    main()
